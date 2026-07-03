#!/usr/bin/env python3.12
"""Build the "JTC Group — Target AI Architecture" deck in four formats from one
content model: docs/jtcgroup_target_ai_architecture.{md,html,pdf,pptx}.

Two reference architectures for production:
  1. Open-source & portable — Container Apps + Blob + PostgreSQL + Azure AI
     Foundry (pluggable LLMs, Anthropic Claude default) + LangGraph + Langfuse.
  2. Microsoft-managed — Foundry Agent Service + Microsoft Fabric + Power BI.

JTC-branded (purple #550055 / #6b1766 / #ba2a84). The HTML deck is 1280×720
landscape slides; the PDF is those slides screenshotted in headless Chromium and
stitched one-per-page; the PPTX is native python-pptx with the same branding.

Usage:
  python3.12 scripts/generate_azure_options.py          # md + html + pdf + pptx
  python3.12 scripts/generate_azure_options.py html     # md + html only
"""
from __future__ import annotations

import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
DOCS = ROOT / "docs"
STEM = "jtcgroup_target_ai_architecture"
VERSION = "v2.0"

# JTC brand palette.
NAVY = "#550055"
NAVY2 = "#6b1766"
ACCENT = "#ba2a84"
INK = "#48484f"
MUTED = "#7a7a85"
GREEN = "#1c7c44"
RED = "#b0353a"

# ── Content model ────────────────────────────────────────────────────────────
# Each slide is a dict: {"kind": ..., "title": ..., ...}. Renderers below map
# each kind to HTML, Markdown, and PPTX.
SLIDES = [
    {"kind": "title", "title": "JTC Group — Target AI Architecture",
     "big2": "Target AI Architecture",
     "subtitle": "Skill-driven AI platform — two reference deployment architectures on Azure",
     "tagline": "Open-source & portable  ·  vs  ·  Microsoft-managed",
     "footer": f"Reference architecture · {VERSION} · Draft for review"},

    {"kind": "bullets", "title": "The baseline platform (both options share this)",
     "intro": "Both architectures run the same skill-driven AI platform — only the "
              "managed services underneath differ.",
     "bullets": [
        "Skill-driven: domain skill packs (single family office lead, tax, "
        "project management) loaded on demand — new domains ship as skills, not apps.",
        "Tools & data reached through MCP servers — one open, model-agnostic "
        "interface, reusable across every skill.",
        "LangGraph orchestrator routes each request to the right skill + agents.",
        "LLM via an OpenAI-compatible layer — the provider is a config switch, "
        "not a code change (Anthropic Claude default).",
        "Containerised app (Docker) with per-team isolation, RBAC, invites, chat "
        "logging, 👍/👎 feedback and LLM-judge evals built in."]},

    {"kind": "components", "title": "Skill library — modular domain skill packs",
     "intro": "The platform is skill-driven: each capability is a versioned skill "
              "pack (instructions + tools + MCP bindings) the orchestrator loads on "
              "demand. New domains ship as new skill packs — not new applications.",
     "rows": [
        ("Single family office lead skills", "Lead sourcing & qualification, "
         "relationship intelligence, prospect research, suitability & mandate "
         "drafting, onboarding / KYC"),
        ("Tax skills", "Form finding, obligation determination, FATCA/CRS "
         "readiness, W-8 preparation, regulatory-change monitoring with citations"),
        ("Project management skills", "Workstream & task planning, deadline / "
         "milestone tracking, status & RAID reporting, resourcing"),
        ("Skill anatomy", "Each pack = a SKILL.md (instructions) + tools + MCP "
         "bindings — hot-swappable, versioned and permission-scoped per team")]},

    {"kind": "components", "title": "MCP servers — standard tool & data access",
     "intro": "Agents reach tools and data through Model Context Protocol (MCP) "
              "servers — a standard, pluggable interface any MCP-capable model "
              "(Anthropic Claude natively) can call. Skills bind to the servers they need.",
     "rows": [
        ("Portfolio / entity MCP", "Entities, obligations, filings, AEOI readiness"),
        ("Document & RAG MCP", "Tax-law corpus, retrieval, citations"),
        ("CRM / relationship MCP", "Contacts, mandates, pipeline — for the family-"
         "office lead skills"),
        ("Filings & authority MCP", "Form catalogues, e-file portals, deadlines"),
        ("Office MCP", "Email, calendar & documents (Microsoft 365 / Google)"),
        ("Market & reference MCP", "External market and reference data")]},

    {"kind": "components", "title": "Option 1 — Open-source & portable",
     "badge": "RECOMMENDED",
     "intro": "Managed Azure PaaS for the plumbing, open-source for the brain — "
              "no lock-in, fully pluggable LLMs with Anthropic Claude as default.",
     "rows": [
        ("Azure Container Apps", "Runs the FastHTML app (same Docker image)"),
        ("Azure Blob Storage", "Form PDFs & scraped source documents"),
        ("Azure Database for PostgreSQL", "Entities, obligations, chat logs, "
         "feedback & evals · pgvector for embeddings · optional Apache AGE for "
         "the citation graph"),
        ("Azure AI Foundry (model catalog)", "Pluggable LLMs — Anthropic Claude "
         "default (Opus 4.8 reasoning / Sonnet 4.6 cost); swap to OpenAI, Llama, "
         "Mistral with no code change"),
        ("LangGraph + skill packs", "Agent orchestrator loads the SFO-lead / tax "
         "/ project-management skill packs on demand"),
        ("MCP servers (containers)", "Standard tool & data access — portfolio, "
         "docs, CRM, filings, office; reusable across every skill"),
        ("Langfuse", "Open-source LLM observability: traces, evals, cost, user "
         "feedback (self-host or Langfuse Cloud)"),
        ("Key Vault + Entra ID", "Secrets & single sign-on")]},

    {"kind": "bullets", "title": "Option 1 — how it flows",
     "intro": "",
     "bullets": [
        "User → Container Apps (FastHTML) over HTTPS with Entra SSO.",
        "LangGraph routes each request to specialist agents; every step is "
        "traced to Langfuse.",
        "LLM calls hit Azure AI Foundry via the OpenAI-compatible endpoint — "
        "Anthropic Claude by default, swappable per-agent.",
        "Retrieval runs over PostgreSQL (pgvector hybrid) + PDFs in Blob.",
        "👍/👎 feedback and LLM-judge scores flow to Langfuse and the in-app "
        "analytics for a full quality loop."]},

    {"kind": "proscons", "title": "Option 1 — pros & cons",
     "pros": [
        "No lock-in — portable across Azure, other clouds, or on-prem.",
        "Truly pluggable LLMs: Anthropic default, swap freely; no single-vendor "
        "model risk.",
        "Reuses the existing codebase (LangGraph, pluggable storage) — minimal "
        "rework.",
        "Open-source observability you own (Langfuse) — full trace/eval/cost data.",
        "Transparent usage-based cost; scale each component independently."],
     "cons": [
        "More moving parts to operate (Container Apps, Postgres, Langfuse).",
        "You own upgrades, scaling & backups — eased by managed PaaS.",
        "Graph queries need pgvector/Apache AGE, or a managed Neo4j AuraDB add-on.",
        "Langfuse is another service to run (or pay for Langfuse Cloud)."]},

    {"kind": "components", "title": "Option 2 — Microsoft-managed",
     "badge": "MANAGED",
     "intro": "Foundry Agent Service + Microsoft Fabric — least ops, deepest "
              "Microsoft integration, at the cost of portability.",
     "rows": [
        ("Azure AI Foundry Agent Service", "Managed agents, built-in tool-calling, "
         "threads & multi-agent orchestration (replaces custom LangGraph)"),
        ("Skills as connected agents", "SFO-lead / tax / project-management skill "
         "packs map to Foundry connected agents & tool sets"),
        ("MCP tools in Foundry", "Foundry Agent Service consumes the same MCP "
         "servers as managed tools"),
        ("Azure AI Search", "Managed vector + hybrid RAG over the corpus"),
        ("Microsoft Fabric (OneLake)", "Unified data lake for tax docs · "
         "Lakehouse / Warehouse · Data Factory ingestion pipelines"),
        ("Power BI", "Firmwide compliance & filing dashboards on Fabric data"),
        ("Azure OpenAI (GPT)", "Default models; Anthropic Claude also available "
         "in the Foundry catalog"),
        ("Azure Monitor + Foundry evaluations", "Built-in tracing & evaluation"),
        ("Microsoft Purview", "Governance, data lineage & DLP")]},

    {"kind": "proscons", "title": "Option 2 — pros & cons",
     "pros": [
        "Fully managed — least ops; Microsoft runs the agents, RAG and scaling.",
        "Deep Microsoft integration: Entra, Purview governance, Fabric, Power BI.",
        "Enterprise SLA & support — strong fit for a Microsoft-shop like JTC.",
        "Built-in agent orchestration, managed RAG and evaluations out of the box.",
        "Fabric unifies data + BI for firmwide reporting."],
     "cons": [
        "Vendor lock-in — Azure/Foundry/Fabric-specific; low portability.",
        "Rework: migrate off LangGraph/Neo4j to Foundry Agents + Fabric.",
        "Model choice steered to OpenAI; less control over orchestration internals.",
        "Fabric capacity (F-SKUs) + Search + Foundry can be costly & hard to "
        "predict.",
        "Managed-agent surface is newer and still maturing."]},

    {"kind": "table", "title": "Side by side",
     "headers": ["Dimension", "Option 1 · Open-source", "Option 2 · MS-managed"],
     "rows": [
        ["Compute", "Container Apps (Docker)", "Foundry Agent Service (managed)"],
        ["Orchestration", "LangGraph (code)", "Foundry Agents (managed)"],
        ["Skills", "Portable skill packs (SKILL.md + tools)",
         "Foundry connected agents / tool sets"],
        ["Tool access", "MCP servers (open standard, reusable)",
         "MCP tools within Foundry"],
        ["LLMs", "Foundry catalog — Anthropic default, fully swappable",
         "Azure OpenAI default; Claude via catalog"],
        ["Data", "PostgreSQL (+pgvector/AGE) + Blob", "Fabric OneLake + AI Search"],
        ["Observability", "Langfuse (open-source)", "Azure Monitor + Foundry evals"],
        ["Analytics / BI", "In-app + Langfuse", "Power BI on Fabric"],
        ["Lock-in", "Low — portable", "High — Azure-native"],
        ["Ops burden", "Higher", "Lower"],
        ["Rework", "Minimal (reuses code)", "Significant"],
        ["Best for", "Portability, model choice, cost control",
         "MS-shop, managed ops, firmwide BI"]]},

    {"kind": "bullets", "title": "Recommendation",
     "intro": "",
     "bullets": [
        "Default to Option 1 (open-source & portable): keeps Anthropic Claude as "
        "the primary model with freedom to swap, reuses the current LangGraph "
        "app, and avoids lock-in. Langfuse adds production-grade observability "
        "on top of the in-app analytics.",
        "Choose Option 2 when JTC wants a fully-managed, Microsoft-native stack "
        "with Fabric / Power BI for firmwide reporting and accepts Azure lock-in "
        "plus the migration effort.",
        "Hybrid path: start on Option 1 (fast, low-risk, portable), then adopt "
        "Fabric / Power BI for analytics later if firmwide BI becomes a priority "
        "— the pluggable storage and OpenAI-compatible LLM layer make that "
        "incremental, not a rebuild.",
        "The domain IP — the family-office lead, tax and project-management skill "
        "packs plus the MCP servers — is portable across both options, so the "
        "investment in skills carries over regardless of the deployment choice."]},
]

# ── HTML deck ────────────────────────────────────────────────────────────────
CSS = f"""
*{{margin:0;padding:0;box-sizing:border-box}}
body{{font-family:-apple-system,'Segoe UI',Helvetica,Arial,sans-serif;background:#33203a}}
.slide{{width:1280px;height:720px;background:#fbfcfd;position:relative;overflow:hidden;
  margin:0 auto 28px;display:flex;flex-direction:column}}
.bar{{height:14px;background:linear-gradient(90deg,{NAVY},{ACCENT})}}
.foot{{position:absolute;bottom:0;left:0;right:0;height:40px;background:{NAVY};color:#fff;
  display:flex;align-items:center;padding:0 40px;font-size:12.5px;letter-spacing:.04em}}
.foot .b{{font-weight:700}}.foot .b span{{color:#ec9bd0}}
.body{{flex:1;padding:30px 44px 52px;display:flex;flex-direction:column}}
h2{{color:{NAVY2};font-size:31px;margin-bottom:6px;display:flex;align-items:center;gap:12px}}
.badge{{font-size:12px;font-weight:700;color:#fff;background:{ACCENT};border-radius:20px;
  padding:3px 12px;letter-spacing:.05em}}
.intro{{color:{INK};font-size:16px;line-height:1.45;margin:4px 0 14px;max-width:1050px}}
.bul li{{list-style:none;color:{INK};font-size:18px;line-height:1.45;margin:11px 0;
  padding-left:24px;position:relative}}
.bul li:before{{content:'';position:absolute;left:0;top:9px;width:10px;height:10px;
  border-radius:50%;background:{ACCENT}}}
.kv{{display:flex;flex-direction:column;gap:8px;margin-top:4px}}
.kv .row{{display:flex;gap:14px;align-items:baseline}}
.kv .svc{{flex:0 0 300px;font-weight:700;color:{NAVY};font-size:16px}}
.kv .role{{flex:1;color:{INK};font-size:15px;line-height:1.35}}
.pc{{display:flex;gap:26px;flex:1;margin-top:6px}}
.pc .col{{flex:1;background:#fff;border:1px solid #e6e3ec;border-radius:12px;padding:16px 20px}}
.pc h3{{font-size:18px;margin-bottom:8px}}
.pc .pros h3{{color:{GREEN}}}.pc .cons h3{{color:{RED}}}
.pc li{{list-style:none;font-size:15px;line-height:1.4;color:{INK};margin:9px 0;padding-left:22px;position:relative}}
.pc .pros li:before{{content:'✓';position:absolute;left:0;color:{GREEN};font-weight:700}}
.pc .cons li:before{{content:'✕';position:absolute;left:0;color:{RED};font-weight:700}}
table{{border-collapse:collapse;width:100%;margin-top:8px;font-size:14.5px}}
th,td{{text-align:left;padding:8px 12px;border-bottom:1px solid #e6e3ec;color:{INK};vertical-align:top}}
th{{background:{NAVY};color:#fff;font-weight:600}}
tr:nth-child(even) td{{background:#f7f4f9}}
td:first-child{{font-weight:700;color:{NAVY2}}}
.title-slide .body{{justify-content:center;padding-left:80px}}
.title-slide .big{{font-size:84px;font-weight:800;color:{NAVY};line-height:1}}
.title-slide .big span{{color:{ACCENT}}}
.title-slide .big2{{font-size:44px;font-weight:800;color:{ACCENT};line-height:1.05;margin-top:6px}}
.title-slide .sub{{font-size:22px;color:{INK};margin-top:18px;max-width:960px}}
.title-slide .tag{{font-size:20px;color:{NAVY2};font-weight:700;margin-top:22px}}
.title-slide .prep{{font-size:16px;color:{MUTED};margin-top:8px}}
"""

FOOT = ('<div class="foot"><span class="b">JTC <span>Group</span></span>'
        f'&nbsp;&nbsp;·&nbsp;&nbsp;Target AI Architecture&nbsp;&nbsp;·&nbsp;&nbsp;{VERSION}</div>')


def _slide_html(s):
    k = s["kind"]
    if k == "title":
        return (f'<div class="slide title-slide"><div class="bar"></div><div class="body">'
                f'<div class="big">JTC <span>Group</span></div>'
                f'<div class="big2">{s["big2"]}</div>'
                f'<div class="sub">{s["subtitle"]}</div>'
                f'<div class="tag">{s["tagline"]}</div>'
                f'<div class="prep">{s["footer"]}</div></div>{FOOT}</div>')
    badge = f'<span class="badge">{s["badge"]}</span>' if s.get("badge") else ""
    head = f'<h2>{s["title"]}{badge}</h2>'
    intro = f'<div class="intro">{s["intro"]}</div>' if s.get("intro") else ""
    if k == "bullets":
        body = '<ul class="bul">' + "".join(f"<li>{b}</li>" for b in s["bullets"]) + "</ul>"
    elif k == "components":
        rows = "".join(f'<div class="row"><div class="svc">{svc}</div>'
                       f'<div class="role">{role}</div></div>' for svc, role in s["rows"])
        body = f'<div class="kv">{rows}</div>'
    elif k == "proscons":
        pros = "".join(f"<li>{p}</li>" for p in s["pros"])
        cons = "".join(f"<li>{c}</li>" for c in s["cons"])
        body = (f'<div class="pc"><div class="col pros"><h3>Pros</h3><ul>{pros}</ul></div>'
                f'<div class="col cons"><h3>Cons</h3><ul>{cons}</ul></div></div>')
    elif k == "table":
        th = "".join(f"<th>{h}</th>" for h in s["headers"])
        trs = "".join("<tr>" + "".join(f"<td>{c}</td>" for c in row) + "</tr>" for row in s["rows"])
        body = f"<table><tr>{th}</tr>{trs}</table>"
    else:
        body = ""
    return (f'<div class="slide"><div class="bar"></div><div class="body">'
            f'{head}{intro}{body}</div>{FOOT}</div>')


def write_html():
    body = "\n".join(_slide_html(s) for s in SLIDES)
    html = (f"<!doctype html><html><head><meta charset='utf-8'>"
            f"<title>JTC Group — Target AI Architecture</title><style>{CSS}</style>"
            f"</head><body>{body}</body></html>")
    (DOCS / f"{STEM}.html").write_text(html)
    print(f"wrote docs/{STEM}.html ({len(SLIDES)} slides)")


# ── Markdown ─────────────────────────────────────────────────────────────────
def write_md():
    out = ["# JTC Group — Target AI Architecture\n",
           "_Two reference architectures for the production AI platform on Azure: "
           "**open-source & portable** vs **Microsoft-managed**._\n",
           f"A slide version is at [`{STEM}.pdf`]({STEM}.pdf) / [`.pptx`]({STEM}.pptx).\n"]
    for s in SLIDES:
        if s["kind"] == "title":
            continue
        badge = f" — {s['badge']}" if s.get("badge") else ""
        out.append(f"## {s['title']}{badge}\n")
        if s.get("intro"):
            out.append(s["intro"] + "\n")
        if s["kind"] == "bullets":
            out += [f"- {b}" for b in s["bullets"]]
        elif s["kind"] == "components":
            out += [f"- **{svc}** — {role}" for svc, role in s["rows"]]
        elif s["kind"] == "proscons":
            out.append("**Pros**\n")
            out += [f"- ✓ {p}" for p in s["pros"]]
            out.append("\n**Cons**\n")
            out += [f"- ✕ {c}" for c in s["cons"]]
        elif s["kind"] == "table":
            out.append("| " + " | ".join(s["headers"]) + " |")
            out.append("|" + "|".join(["---"] * len(s["headers"])) + "|")
            out += ["| " + " | ".join(r) + " |" for r in s["rows"]]
        out.append("")
    (DOCS / f"{STEM}.md").write_text("\n".join(out))
    print(f"wrote docs/{STEM}.md")


# ── PDF (screenshot each HTML slide, stitch landscape) ───────────────────────
def write_pdf():
    from fpdf import FPDF
    from playwright.sync_api import sync_playwright
    import io

    pngs = []
    html_uri = (DOCS / f"{STEM}.html").as_uri()
    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(viewport={"width": 1280, "height": 720}, device_scale_factor=2)
        page.goto(html_uri, wait_until="networkidle")
        slides = page.locator(".slide")
        for i in range(slides.count()):
            pngs.append(slides.nth(i).screenshot())
            print(f"  rendered slide {i + 1}/{slides.count()}")
        browser.close()

    pw, ph = 338.667, 190.5  # 16:9 landscape mm
    pdf = FPDF(orientation="L", unit="mm", format=(ph, pw))
    pdf.set_auto_page_break(False)
    for png in pngs:
        pdf.add_page()
        pdf.image(io.BytesIO(png), x=0, y=0, w=pw, h=ph)
    pdf.output(str(DOCS / f"{STEM}.pdf"))
    print(f"wrote docs/{STEM}.pdf ({len(pngs)} slides)")


# ── PPTX (native, JTC-branded) ───────────────────────────────────────────────
def _rgb(hexstr):
    from pptx.dml.color import RGBColor
    return RGBColor.from_string(hexstr.lstrip("#").upper())


def write_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    def bar_and_foot(slide):
        bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(0.16))
        bar.fill.solid(); bar.fill.fore_color.rgb = _rgb(NAVY)
        bar.line.fill.background(); bar.shadow.inherit = False
        foot = slide.shapes.add_shape(1, 0, SH - Inches(0.34), SW, Inches(0.34))
        foot.fill.solid(); foot.fill.fore_color.rgb = _rgb(NAVY)
        foot.line.fill.background(); foot.shadow.inherit = False
        tf = foot.text_frame; tf.margin_top = Pt(2); tf.margin_left = Inches(0.4)
        r = tf.paragraphs[0].add_run()
        r.text = f"JTC Group · Target AI Architecture · {VERSION}"
        r.font.size = Pt(9); r.font.color.rgb = _rgb("FFFFFF")

    def add_title(slide, text, badge=None):
        box = slide.shapes.add_textbox(Inches(0.5), Inches(0.34), Inches(12.3), Inches(0.9))
        p = box.text_frame.paragraphs[0]
        r = p.add_run(); r.text = text
        r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = _rgb(NAVY2)
        if badge:
            rb = p.add_run(); rb.text = "   " + badge
            rb.font.size = Pt(13); rb.font.bold = True; rb.font.color.rgb = _rgb(ACCENT)

    def body_frame(slide, top=Inches(1.3), height=Inches(5.6)):
        box = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.3), height)
        tf = box.text_frame; tf.word_wrap = True
        return tf

    for s in SLIDES:
        slide = prs.slides.add_slide(blank)
        bar_and_foot(slide)
        k = s["kind"]

        if k == "title":
            box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(3))
            tf = box.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = "JTC Group"
            r.font.size = Pt(54); r.font.bold = True; r.font.color.rgb = _rgb(NAVY)
            for txt, sz, col, bold in [
                (s["big2"], 34, ACCENT, True),
                (s["subtitle"], 20, INK, False),
                (s["tagline"], 18, NAVY2, True),
                (s["footer"], 14, MUTED, False)]:
                pp = tf.add_paragraph(); rr = pp.add_run(); rr.text = txt
                rr.font.size = Pt(sz); rr.font.color.rgb = _rgb(col); rr.font.bold = bold
                pp.space_before = Pt(10)
            continue

        add_title(slide, s["title"], s.get("badge"))
        tf = body_frame(slide)
        first = True

        def para(text, *, size=16, color=INK, bold=False, bullet=False, sym=None,
                 symcolor=None, space=8):
            nonlocal first
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = Pt(space)
            if sym:
                rs = p.add_run(); rs.text = sym + " "
                rs.font.size = Pt(size); rs.font.bold = True
                rs.font.color.rgb = _rgb(symcolor or ACCENT)
            elif bullet:
                rs = p.add_run(); rs.text = "•  "
                rs.font.size = Pt(size); rs.font.color.rgb = _rgb(ACCENT)
            r = p.add_run(); r.text = text
            r.font.size = Pt(size); r.font.color.rgb = _rgb(color); r.font.bold = bold
            return p

        if s.get("intro"):
            para(s["intro"], size=15, color=INK, space=12)

        if k == "bullets":
            for b in s["bullets"]:
                para(b, size=16, bullet=True, space=9)
        elif k == "components":
            for svc, role in s["rows"]:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.space_after = Pt(7)
                rs = p.add_run(); rs.text = svc + "  —  "
                rs.font.size = Pt(15); rs.font.bold = True; rs.font.color.rgb = _rgb(NAVY)
                rr = p.add_run(); rr.text = role
                rr.font.size = Pt(14); rr.font.color.rgb = _rgb(INK)
        elif k == "proscons":
            # two side-by-side text boxes
            for i, (label, items, col, sym) in enumerate([
                    ("Pros", s["pros"], GREEN, "✓"), ("Cons", s["cons"], RED, "✕")]):
                bx = slide.shapes.add_textbox(Inches(0.5 + i * 6.35), Inches(1.4),
                                              Inches(6.05), Inches(5.4))
                b = bx.text_frame; b.word_wrap = True
                hp = b.paragraphs[0]; hr = hp.add_run(); hr.text = label
                hr.font.size = Pt(19); hr.font.bold = True; hr.font.color.rgb = _rgb(col)
                hp.space_after = Pt(8)
                for it in items:
                    pp = b.add_paragraph(); pp.space_after = Pt(7)
                    rs = pp.add_run(); rs.text = sym + "  "
                    rs.font.size = Pt(14); rs.font.bold = True; rs.font.color.rgb = _rgb(col)
                    rr = pp.add_run(); rr.text = it
                    rr.font.size = Pt(14); rr.font.color.rgb = _rgb(INK)
        elif k == "table":
            rows, cols = len(s["rows"]) + 1, len(s["headers"])
            gt = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.35),
                                        Inches(12.3), Inches(5.4)).table
            gt.columns[0].width = Inches(2.5)
            gt.columns[1].width = Inches(5.1)
            gt.columns[2].width = Inches(4.7)
            for j, h in enumerate(s["headers"]):
                cell = gt.cell(0, j)
                cell.fill.solid(); cell.fill.fore_color.rgb = _rgb(NAVY)
                pr = cell.text_frame.paragraphs[0]; run = pr.add_run(); run.text = h
                run.font.size = Pt(13); run.font.bold = True; run.font.color.rgb = _rgb("FFFFFF")
            for i, row in enumerate(s["rows"], start=1):
                for j, val in enumerate(row):
                    cell = gt.cell(i, j)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _rgb("F7F4F9" if i % 2 else "FFFFFF")
                    pr = cell.text_frame.paragraphs[0]; run = pr.add_run(); run.text = val
                    run.font.size = Pt(11.5)
                    run.font.color.rgb = _rgb(NAVY2 if j == 0 else INK)
                    run.font.bold = (j == 0)

    prs.save(str(DOCS / f"{STEM}.pptx"))
    print(f"wrote docs/{STEM}.pptx ({len(SLIDES)} slides)")


if __name__ == "__main__":
    cmd = sys.argv[1] if len(sys.argv) > 1 else "all"
    write_html()
    write_md()
    if cmd != "html":
        write_pdf()
        write_pptx()
