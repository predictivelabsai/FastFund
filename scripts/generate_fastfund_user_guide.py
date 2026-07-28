#!/usr/bin/env python3
"""Generate the unified FastFund user guide in MD, HTML, PDF and PPTX."""

from __future__ import annotations

import html
import shutil
import subprocess
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
DOCS = ROOT / "docs"
SHOTS = DOCS / "screenshots"

NAVY = "102A43"
NAVY_DARK = "0B1F33"
TEAL = "0F766E"
TEXT = "34495E"
MUTED = "6B7C8F"
BG = "F7FAFC"

# id, title, screenshot, bullets
SLIDES = [
    ("title", "FastFund", None, [
        "Family-office outreach and multijurisdiction tax filing intelligence",
        "One open-source product, one relationship-to-filing data model",
        "User Guide",
    ]),
    ("journey", "One connected operating journey", None, [
        "Onboard a family office and capture principals, portfolio context and service needs.",
        "Develop explainable outreach opportunities, proposals and relationship actions.",
        "Link every fund, trust, company, SPV and holding vehicle to its owning family.",
        "Determine jurisdiction-specific obligations, deadlines and filing readiness.",
        "Prepare, review, file and verify while monitoring changes in source law.",
    ]),
    ("assistant", "FastFund assistant", "fastfund-assistant.png", [
        "A shared assistant routes questions to tax-law, form, change and data specialists.",
        "The rationalised sidebar groups work into Family Office, Tax & Compliance, and Knowledge & Insights.",
        "Answers remain grounded in official source material and linked operational data.",
    ]),
    ("family-advisor", "Family-office advisor", "fastfund-family-advisor.png", [
        "Ask about a family's governance, service gaps, portfolio or next best action.",
        "Rules and the cross-sell graph produce traceable recommendations.",
        "The profile and live recommendations stay visible beside the conversation.",
    ]),
    ("families", "Families and outreach book", "fastfund-families.png", [
        "Manage leads, onboarding relationships and established family-office clients.",
        "Review AUM, domicile, generations, family size and services already held.",
        "Open any family directly in the advisor or its complete relationship profile.",
    ]),
    ("profile", "Family profile and ownership", "fastfund-family-profile.png", [
        "Portfolio allocation, holdings, transactions, members and relationship history in one view.",
        "Outreach recommendations and scheduled actions share the same family context.",
        "Linked legal entities connect relationship work directly to tax filings.",
    ]),
    ("pipeline", "Outreach pipeline", "fastfund-outreach-pipeline.png", [
        "Move recommendations through suggested, presented, accepted and booked stages.",
        "Filter cross-sell and upsell opportunities by service category.",
        "Generate proposals and schedule consultations or follow-ups.",
    ]),
    ("entities", "Legal entities", "fastfund-entities.png", [
        "Funds, trusts, companies, GPs, SPVs and holding vehicles share one portfolio.",
        "Each entity can be linked to its owning family office.",
        "Jurisdictions, activities and financial year-end drive deterministic filing obligations.",
    ]),
    ("obligations", "Filing obligations", "fastfund-obligations.png", [
        "Track every filing from determination through preparation, filing and confirmation.",
        "Statuses and expert verification survive repeat determination runs.",
        "Filter by jurisdiction, category, urgency, entity and completion state.",
    ]),
    ("calendar", "Multijurisdiction filing calendar", "fastfund-filing-calendar.png", [
        "Resolve form rules into real deadlines using the entity's financial year-end.",
        "Urgency indicators surface overdue, due-soon and upcoming work.",
        "Calendar and register views provide operational and portfolio-level perspectives.",
    ]),
    ("dashboard", "Tax and compliance dashboard", "fastfund-tax-dashboard.png", [
        "See entities, open obligations, filed work and verification coverage at a glance.",
        "Move from portfolio metrics directly to the records behind them.",
        "Use audit, team and assistant analytics for controlled operations.",
    ]),
    ("knowledge", "Knowledge, provenance and regulatory change", None, [
        "Official forms, legislation, guidance, gazettes and treaties are captured by jurisdiction.",
        "Immutable versions and change records preserve the regulatory audit trail.",
        "Citations and provenance edges connect answers and forms back to source law.",
        "SQLite supports zero-infrastructure use; Neo4j adds native graph traversal.",
    ]),
    ("finish", "Human-controlled automation", None, [
        "FastFund assists with research, determination, preparation and outreach.",
        "Users review recommendations, validate form readiness and confirm filing status.",
        "Synthetic demo data is used for family-office examples; official tax sources retain provenance.",
    ]),
]


def _clean(text: str) -> str:
    return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def write_markdown() -> None:
    lines = [
        "# FastFund User Guide",
        "",
        "FastFund combines family-office relationship outreach with multijurisdiction "
        "tax filing intelligence in one open-source product.",
        "",
    ]
    for sid, title, image, bullets in SLIDES[1:]:
        lines.extend([f"## {title}", ""])
        if image:
            lines.extend([f"![{title}](screenshots/{image})", ""])
        lines.extend(f"- {item}" for item in bullets)
        lines.append("")
    (DOCS / "fastfund_user_guide.md").write_text("\n".join(lines), encoding="utf-8")


def _slide_html(sid: str, title: str, image: str | None, bullets: list[str]) -> str:
    items = "".join(f"<li>{html.escape(item)}</li>" for item in bullets)
    image_html = (
        f'<div class="shot"><img src="screenshots/{image}" alt="{html.escape(title)}"></div>'
        if image else ""
    )
    cls = "slide title-slide" if sid == "title" else "slide"
    title_html = (
        '<div class="product">Fast<span>Fund</span></div>'
        if sid == "title" else f"<h2>{html.escape(title)}</h2>"
    )
    return (
        f'<section class="{cls}" id="{sid}"><div class="accent"></div>'
        f'<div class="body">{title_html}<div class="split"><ul>{items}</ul>{image_html}</div></div>'
        '<footer><b>Fast<span>Fund</span></b><i>User Guide</i></footer></section>'
    )


def write_html() -> None:
    slides = "\n".join(_slide_html(*slide) for slide in SLIDES)
    css = f"""
*{{box-sizing:border-box}}html,body{{margin:0;background:#081b2b;font-family:Arial,sans-serif}}
.slide{{width:1280px;height:720px;margin:0 auto 24px;background:#{BG};position:relative;
overflow:hidden;break-after:page;page-break-after:always}}.accent{{height:14px;
background:linear-gradient(90deg,#{NAVY_DARK},#{TEAL})}}.body{{padding:38px 48px 62px}}
h2{{font-size:34px;color:#{NAVY};margin:0 0 22px}}.split{{display:flex;gap:30px;align-items:center}}
ul{{flex:0 0 39%;margin:0;padding:0;list-style:none}}li{{font-size:18px;color:#{TEXT};
line-height:1.45;margin:15px 0;padding-left:20px;position:relative}}li:before{{content:'';
position:absolute;left:0;top:8px;width:9px;height:9px;border-radius:50%;background:#{TEAL}}}
.shot{{flex:1;text-align:center}}.shot img{{max-width:100%;max-height:545px;border:1px solid #d8e2ea;
border-radius:10px;box-shadow:0 8px 24px rgba(16,42,67,.16)}}footer{{position:absolute;left:0;
right:0;bottom:0;height:42px;background:#{NAVY_DARK};color:white;padding:0 30px;
display:flex;align-items:center;justify-content:space-between}}footer span,.product span{{color:#{TEAL}}}
footer i{{font-style:normal;color:#b9d8d3}}.title-slide .body{{padding:120px 78px}}
.product{{font-size:92px;line-height:1;font-weight:800;color:#{NAVY_DARK};margin-bottom:30px}}
.title-slide ul{{flex-basis:75%}}.title-slide li{{font-size:25px;margin:20px 0}}
@page{{size:1280px 720px;margin:0}}@media print{{html,body{{background:white}}.slide{{margin:0}}}}
"""
    out = (
        "<!doctype html><html><head><meta charset='utf-8'>"
        f"<title>FastFund User Guide</title><style>{css}</style></head><body>{slides}</body></html>"
    )
    (DOCS / "fastfund_user_guide.html").write_text(out, encoding="utf-8")


def write_pdf() -> None:
    browser = next(
        (name for name in ("chromium", "google-chrome", "chromium-browser")
         if shutil.which(name)),
        None,
    )
    if not browser:
        raise RuntimeError("Chromium or Chrome is required to generate the PDF")
    source = (DOCS / "fastfund_user_guide.html").resolve().as_uri()
    output = DOCS / "fastfund_user_guide.pdf"
    subprocess.run(
        [
            browser,
            "--headless=new",
            "--no-sandbox",
            "--disable-gpu",
            "--no-pdf-header-footer",
            f"--print-to-pdf={output}",
            source,
        ],
        check=True,
    )


def write_powerpoint() -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for sid, title, image, bullets in SLIDES:
        slide = prs.slides.add_slide(blank)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor.from_string(BG)
        bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(.14))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor.from_string(TEAL)
        bar.line.fill.background()
        title_box = slide.shapes.add_textbox(Inches(.55), Inches(.38), Inches(12.2), Inches(.55))
        p = title_box.text_frame.paragraphs[0]
        p.text = "FastFund" if sid == "title" else title
        p.font.name = "Aptos Display"
        p.font.bold = True
        p.font.size = Pt(48 if sid == "title" else 28)
        p.font.color.rgb = RGBColor.from_string(NAVY)
        body_width = Inches(5.0 if image else 11.9)
        box = slide.shapes.add_textbox(Inches(.62), Inches(1.25), body_width, Inches(5.45))
        tf = box.text_frame
        tf.word_wrap = True
        tf.clear()
        for index, item in enumerate(bullets):
            para = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
            para.text = f"• {item}"
            para.level = 0
            para.font.name = "Aptos"
            para.font.size = Pt(18 if sid != "title" else 24)
            para.font.color.rgb = RGBColor.from_string(TEXT)
            para.space_after = Pt(12)
        if image:
            path = SHOTS / image
            if path.exists():
                slide.shapes.add_picture(str(path), Inches(5.8), Inches(1.12),
                                         width=Inches(7.0), height=Inches(5.55))
        footer = slide.shapes.add_textbox(Inches(.55), Inches(7.05), Inches(12.2), Inches(.25))
        fp = footer.text_frame.paragraphs[0]
        fp.text = "FastFund · User Guide"
        fp.alignment = PP_ALIGN.RIGHT
        fp.font.size = Pt(10)
        fp.font.color.rgb = RGBColor.from_string(MUTED)
    prs.save(DOCS / "fastfund_user_guide.pptx")


def main() -> None:
    selected = set(sys.argv[1:] or ("all",))
    write_markdown()
    write_html()
    if "all" in selected or "pdf" in selected:
        write_pdf()
    if "all" in selected or "pptx" in selected:
        write_powerpoint()
    print("Generated FastFund guide: MD, HTML"
          + (", PDF" if "all" in selected or "pdf" in selected else "")
          + (", PPTX" if "all" in selected or "pptx" in selected else ""))


if __name__ == "__main__":
    main()
